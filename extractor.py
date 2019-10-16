import json
import subprocess as sp
from pathlib import Path
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfFileReader
from utils import clean_line


def process_doc(file_path: Path):
    """ Extracts text from .doc files
    Args:
        file_path(Path) : Path object that contains the file_path of the .doc file
    Returns:
        list : The sentences extracted from the file
    """
    try:

        p = sp.run(["catdoc", str(file_path)], capture_output=True)
        output = p.stdout.decode()
        sentences = [clean_line(line) for line in output.split("\n\n") if line]
        return sentences
    except FileNotFoundError as e:
        print("Unable to process", file_path)
        print(e.strerror)
        return []


def process_docx(file_path: Path):
    """ Extracts text from .docx files
    Args:
        file_path(Path) : Path object that contains the file_path of the .docx file
    Returns:
        list : The sentences extracted from the file
    """
    doc = Document(file_path)
    sentences = []
    for para in doc.paragraphs:
        for line in para.text.split("."):
            line = clean_line(line)
            if line:
                sentences.append(line)
    return sentences


def process_pptx(file_path: Path):
    """ Extracts text from .pptx files
    Args:
        file_path(Path) : Path object that contains the file_path of the .pptx file
    Returns:
        list : The sentences extracted from the file
    """
    prs = Presentation(file_path)
    text_runs = []
    for i, slide in enumerate(prs.slides):
        if i < 1:
            continue
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                para_text = "".join(run.text for run in paragraph.runs)
                if para_text:
                    text_runs.append(para_text)
    return text_runs


def process_pdf(file_path: Path):
    """ Extracts text from .pdf files
    Args:
        file_path(Path) : Path object that contains the file_path of the .pdf file
    Returns:
        list : The sentences extracted from the file
    """
    sentences = []
    with open(file_path, "rb") as f:
        reader = PdfFileReader(f)
        for page in reader.pages:
            for line in page.extractText().splitlines():
                line = clean_line(line)
                if line not in ("", "?", ".", "-", ",", ":"):
                    sentences.append(line)
    return sentences


def extract_sentences(file_path: Path):
    """ Extract sentences from file and store in JSON
    Args:
        file_path(Path) : Path object that contains the path of the file
    Returns:
        list : List of all sentences extracted from the file
    """
    json_file = file_path.with_suffix(".json")
    if json_file.exists():
        with open(json_file) as f:
            return json.load(f)["sentences"]

    if file_path.suffix == ".pptx":
        sentences = process_pptx(file_path)
    elif file_path.suffix == ".doc":
        sentences = process_doc(file_path)
    elif file_path.suffix == ".pdf":
        sentences = process_pdf(file_path)
    if not sentences:
        return []
    data = {"name": str(file_path), "sentences": sentences}

    with open(json_file, "w") as f:
        json.dump(data, f, indent=4)  # Create/Dump json file with file's text
    return sentences

